from __future__ import annotations

from collections import Counter
import hashlib
import os
import re
import tempfile
from pathlib import Path

import fitz
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from classcorpus.models import ExtractionStatus, SlideRecord, VisualAsset
from classcorpus.parser_plugins import TEXT_PLUGIN
from classcorpus.parser_registry import ParserPlugin, ParserRegistry


class UnsupportedFormatError(ValueError):
    def __init__(self, suffix: str):
        label = suffix or "<no extension>"
        super().__init__(f"Unsupported source format: {label}")


def parse_source(path: Path, render_dir: Path) -> list[SlideRecord]:
    plugin = _PARSER_REGISTRY.parser_for(path.suffix)
    if plugin is None:
        raise UnsupportedFormatError(path.suffix.lower())
    return plugin.parse(path, render_dir)


def _parse_pdf(path: Path, render_dir: Path) -> list[SlideRecord]:
    render_dir.mkdir(parents=True, exist_ok=True)

    records: list[SlideRecord] = []
    with fitz.open(path) as document:
        for ordinal, page in enumerate(document, start=1):
            raw_text = page.get_text("text", sort=True)
            page_text = _normalized_text(raw_text)
            title, body_text = _split_title_and_body(page_text)
            has_visual_content, extraction_reasons = _pdf_audit(page, raw_text)

            image_path = render_dir / f"page-{ordinal:04d}.png"
            page.get_pixmap(
                matrix=fitz.Matrix(1.5, 1.5),
                alpha=False,
            ).save(image_path)

            records.append(
                SlideRecord(
                    ordinal=ordinal,
                    kind="page",
                    title=title,
                    body_text=body_text,
                    speaker_notes="",
                    raw_text=raw_text,
                    extraction_status=_status(extraction_reasons),
                    extraction_reasons=extraction_reasons,
                    native_text_chars=len(raw_text),
                    has_visual_content=has_visual_content,
                    render_path=str(image_path),
                )
            )

    return records


def _parse_pptx(path: Path, render_dir: Path) -> list[SlideRecord]:
    presentation = Presentation(path)

    records: list[SlideRecord] = []
    for ordinal, slide in enumerate(presentation.slides, start=1):
        text_frames: list[str] = []
        table_texts: list[str] = []
        extracted_texts: list[str] = []

        for shape in slide.shapes:
            _collect_shape_text(shape, text_frames, table_texts, extracted_texts)

        title = text_frames[0] if text_frames else ""
        body_parts = text_frames[1:] + table_texts
        census = _xml_texts(slide.element)
        missing_texts = list(dict.fromkeys(_missing_texts(census, extracted_texts)))
        body_parts.extend(missing_texts)
        raw_text = "\n".join(census)
        has_visual_content, extraction_reasons = _pptx_audit(
            slide,
            missing_texts,
        )
        visual_assets = _collect_visual_assets(slide.shapes, render_dir)

        notes_text = ""
        try:
            notes_text = _normalized_text(slide.notes_slide.notes_text_frame.text)
        except AttributeError:
            notes_text = ""

        records.append(
            SlideRecord(
                ordinal=ordinal,
                kind="slide",
                title=title,
                body_text="\n".join(body_parts),
                speaker_notes=notes_text,
                raw_text=raw_text,
                extraction_status=_status(extraction_reasons),
                extraction_reasons=extraction_reasons,
                native_text_chars=len(raw_text),
                has_visual_content=has_visual_content,
                render_path=None,
                visual_assets=visual_assets,
            )
        )

    return records


def _collect_shape_text(
    shape,
    text_frames: list[str],
    table_texts: list[str],
    extracted_texts: list[str],
) -> None:
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:
            _collect_shape_text(child, text_frames, table_texts, extracted_texts)
        return

    if getattr(shape, "has_text_frame", False):
        text = _normalized_text(shape.text_frame.text)
        if text:
            text_frames.append(text)
        extracted_texts.extend(_text_frame_texts(shape.text_frame))

    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            for cell in row.cells:
                text = _normalized_text(cell.text)
                if text:
                    table_texts.append(text)
                extracted_texts.extend(_text_frame_texts(cell.text_frame))


def _text_frame_texts(text_frame) -> list[str]:
    texts: list[str] = []
    for paragraph in text_frame.paragraphs:
        run_texts = [run.text for run in paragraph.runs if run.text]
        if run_texts:
            texts.extend(run_texts)
        elif paragraph.text:
            texts.append(paragraph.text)
    return texts


def _tokens(text: str) -> Counter[str]:
    return Counter(re.findall(r"\w+", text.casefold(), flags=re.UNICODE))


def _pdf_audit(page: fitz.Page, raw_text: str) -> tuple[bool, tuple[str, ...]]:
    reasons: list[str] = []
    images = bool(page.get_images(full=True))
    drawings = bool(page.get_drawings())
    word_text = " ".join(str(word[4]) for word in page.get_text("words"))
    block_text = "\n".join(
        str(block[4]) for block in page.get_text("blocks") if len(block) > 4
    )
    native_tokens = _tokens(raw_text)
    if _tokens(word_text) - native_tokens or _tokens(block_text) - native_tokens:
        reasons.append("native-extractor-disagreement")
    if not raw_text.strip():
        reasons.append("no-native-text")
    elif images and len(raw_text.strip()) < 80:
        reasons.append("low-native-text")
        reasons.append("embedded-image")
    return images or drawings, tuple(dict.fromkeys(reasons))


def _xml_texts(element) -> list[str]:
    return [
        node.text
        for node in element.iter()
        if node.tag.endswith("}t") and node.text
    ]


def _missing_texts(census: list[str], extracted: list[str]) -> list[str]:
    remaining = Counter(text.strip() for text in extracted if text.strip())
    missing: list[str] = []
    for text in (item.strip() for item in census):
        if not text:
            continue
        if remaining[text]:
            remaining[text] -= 1
        else:
            missing.append(text)
    return missing


def _pptx_audit(slide, missing_texts: list[str]) -> tuple[bool, tuple[str, ...]]:
    shape_types = set(_shape_types(slide.shapes))
    local_tags = {
        node.tag.rsplit("}", 1)[-1].casefold() for node in slide.element.iter()
    }
    namespaces = {
        node.tag[1:].split("}", 1)[0].casefold()
        for node in slide.element.iter()
        if node.tag.startswith("{")
    }
    relationship_suffixes = {
        relationship.reltype.rsplit("/", 1)[-1].casefold()
        for relationship in slide.part.rels.values()
    }

    has_image = (
        bool(
            shape_types
            & {
                MSO_SHAPE_TYPE.PICTURE,
                MSO_SHAPE_TYPE.LINKED_PICTURE,
            }
        )
        or "pic" in local_tags
        or "image" in relationship_suffixes
    )
    has_chart_or_diagram = (
        bool(
            shape_types
            & {
                MSO_SHAPE_TYPE.CHART,
                MSO_SHAPE_TYPE.DIAGRAM,
                MSO_SHAPE_TYPE.IGX_GRAPHIC,
            }
        )
        or "chart" in local_tags
        or "http://schemas.openxmlformats.org/drawingml/2006/diagram"
        in namespaces
        or bool(
            relationship_suffixes
            & {
                "chart",
                "diagramdata",
                "diagramlayout",
                "diagramcolors",
                "diagramquickstyle",
            }
        )
    )
    has_equation_or_object = (
        bool(
            shape_types
            & {
                MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT,
                MSO_SHAPE_TYPE.LINKED_OLE_OBJECT,
                MSO_SHAPE_TYPE.OLE_CONTROL_OBJECT,
            }
        )
        or bool(local_tags & {"omath", "oleobj"})
        or bool(relationship_suffixes & {"oleobject", "package"})
    )

    reasons: list[str] = []
    if missing_texts:
        reasons.append("unmapped-ooxml-text")
    if has_image:
        reasons.append("embedded-image")
    if has_chart_or_diagram:
        reasons.append("chart-or-diagram")
    if has_equation_or_object:
        reasons.append("equation-or-embedded-object")

    return (
        has_image or has_chart_or_diagram or has_equation_or_object,
        tuple(reasons),
    )


def _shape_types(shapes):
    for shape in shapes:
        yield shape.shape_type
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _shape_types(shape.shapes)


def _collect_visual_assets(
    shapes,
    render_dir: Path,
) -> tuple[VisualAsset, ...]:
    assets: list[VisualAsset] = []
    asset_dir = render_dir / "assets"

    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            assets.extend(_collect_visual_assets(shape.shapes, render_dir))
            continue
        if shape.shape_type not in {
            MSO_SHAPE_TYPE.PICTURE,
            MSO_SHAPE_TYPE.LINKED_PICTURE,
        }:
            continue

        image = shape.image
        blob = image.blob
        digest = hashlib.sha256(blob).hexdigest()
        extension = image.ext or "bin"
        asset_dir.mkdir(parents=True, exist_ok=True)
        asset_path = asset_dir / f"{digest}.{extension}"
        if not asset_path.exists():
            with tempfile.NamedTemporaryFile(
                dir=asset_dir,
                prefix=f".{digest}.",
                suffix=".tmp",
                delete=False,
            ) as stream:
                stream.write(blob)
                temporary_path = Path(stream.name)
            try:
                os.replace(temporary_path, asset_path)
            finally:
                temporary_path.unlink(missing_ok=True)

        assets.append(
            VisualAsset(
                path=str(asset_path),
                kind="image",
                shape_name=str(shape.name),
                content_type=str(image.content_type),
                left=int(shape.left),
                top=int(shape.top),
                width=int(shape.width),
                height=int(shape.height),
            )
        )

    return tuple(assets)


def _status(reasons: tuple[str, ...]) -> ExtractionStatus:
    return "review-needed" if reasons else "text-extracted"


def _split_title_and_body(text: str) -> tuple[str, str]:
    lines = [line.strip() for line in text.splitlines() if line.strip()]
    if not lines:
        return "", ""
    return lines[0], "\n".join(lines[1:])


def _normalized_text(text: str) -> str:
    return "\n".join(line.strip() for line in text.splitlines() if line.strip())


_PARSER_REGISTRY = ParserRegistry()
_PARSER_REGISTRY.register(
    ParserPlugin(name="pdf", suffixes=(".pdf",), parse=_parse_pdf)
)
_PARSER_REGISTRY.register(
    ParserPlugin(name="powerpoint", suffixes=(".pptx",), parse=_parse_pptx)
)
_PARSER_REGISTRY.register(TEXT_PLUGIN)


def register_parser(plugin: ParserPlugin) -> None:
    _PARSER_REGISTRY.register(plugin)


def supported_suffixes() -> frozenset[str]:
    return _PARSER_REGISTRY.supported_suffixes()


__all__ = [
    "UnsupportedFormatError",
    "parse_source",
    "register_parser",
    "supported_suffixes",
]
