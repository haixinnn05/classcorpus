"""Generate and index a synthetic course so ClassCorpus can be evaluated.

The corpus is written from code, contains no third-party material, and needs no
network access or model download. It lives inside the installed package because
the benchmark corpus under ``benchmarks/`` is not distributed.
"""

from __future__ import annotations

from dataclasses import asdict
from pathlib import Path
from typing import Any

import fitz
from pptx import Presentation
from pptx.util import Inches

from classcorpus.database import Database
from classcorpus.indexer import sync_course
from classcorpus.paths import data_root
from classcorpus.payloads import search_response
from classcorpus.search import search

DEMO_COURSE_NAME = "ClassCorpus Demo"
DEMO_QUERY = "why does Dijkstra fail on negative edges"
DEMO_MARKER_NAME = ".classcorpus-demo"
DEMO_CORPUS_VERSION = "1"


def demo_source_root() -> Path:
    """Return the default generated location for the demo course files."""
    return data_root() / "demo-course"


def generate_demo_corpus(
    directory: Path, *, overwrite: bool = False
) -> tuple[str, ...]:
    """Write the demo lecture files and return their names in stable order.

    Refuses to write into an existing non-empty directory that ClassCorpus did
    not generate, unless ``overwrite`` is explicit.
    """
    marker = directory / DEMO_MARKER_NAME
    if directory.exists() and any(directory.iterdir()) and not marker.exists():
        if not overwrite:
            raise ValueError(
                f"refusing to write demo files into non-empty directory: "
                f"{directory}. Choose an empty --dir or pass --overwrite."
            )
    directory.mkdir(parents=True, exist_ok=True)
    marker.write_text(
        f"ClassCorpus demo corpus version {DEMO_CORPUS_VERSION}\n",
        encoding="utf-8",
    )

    _write_pdf(directory / "handout-shortest-paths.pdf")
    _write_pptx(directory / "Lecture01-Complexity.pptx")
    _write_markdown(directory / "study-notes.md")
    return (
        "Lecture01-Complexity.pptx",
        "handout-shortest-paths.pdf",
        "study-notes.md",
    )


def run_demo(
    database: Database,
    *,
    course: str = DEMO_COURSE_NAME,
    source_root: Path | None = None,
    query: str = DEMO_QUERY,
    overwrite: bool = False,
) -> dict[str, Any]:
    """Generate the demo corpus, index it, and search it once."""
    directory = (source_root or demo_source_root()).expanduser().resolve()
    generated = generate_demo_corpus(directory, overwrite=overwrite)
    report = sync_course(database, course, directory)
    results = search(database, query, course=course, limit=3)
    health = database.source_health(course)
    payload: dict[str, Any] = {
        "ok": report.failed == 0 and bool(results),
        "course": course,
        "source_root": str(directory),
        "generated_files": list(generated),
        "sync": asdict(report),
        "query": query,
        "search": search_response(
            results,
            warnings=[],
            sync_required=health.total == 0 or health.failed > 0,
            suggested_terms=[],
        ),
        "next_steps": [
            f'classcorpus search "memoization" --course "{course}"',
            f'classcorpus outline "{course}"',
            f'classcorpus status --course "{course}"',
            f'classcorpus remove "{course}" --confirm',
        ],
    }
    if report.failed:
        payload["error"] = {
            "type": "PartialSyncError",
            "message": f"{report.failed} demo source file(s) failed to index",
        }
    elif not results:
        payload["error"] = {
            "type": "DemoSearchError",
            "message": f"demo corpus indexed but returned no match for: {query}",
        }
    return payload


def _write_pdf(path: Path) -> None:
    document = fitz.open()
    try:
        for heading, body in _PDF_PAGES:
            _add_pdf_page(document, heading, body)
        document.save(path)
    finally:
        document.close()


def _add_pdf_page(document: Any, heading: str, body: str) -> None:
    page = document.new_page(width=612, height=792)
    page.insert_textbox(
        fitz.Rect(72, 64, 540, 104),
        heading,
        fontsize=16,
        fontname="hebo",
    )
    overflow = page.insert_textbox(
        fitz.Rect(72, 112, 540, 720),
        body,
        fontsize=11,
        fontname="helv",
    )
    if overflow < 0:
        raise RuntimeError(f"demo PDF page does not fit its frame: {heading}")


def _write_pptx(path: Path) -> None:
    presentation = Presentation()
    title_only_layout = presentation.slide_layouts[5]

    for title, bullets, note in _PPTX_SLIDES:
        slide = presentation.slides.add_slide(title_only_layout)
        slide.shapes.title.text = title
        textbox = slide.shapes.add_textbox(
            Inches(0.9),
            Inches(1.7),
            Inches(8.2),
            Inches(3.6),
        )
        frame = textbox.text_frame
        frame.word_wrap = True
        frame.text = bullets[0]
        for bullet in bullets[1:]:
            frame.add_paragraph().text = bullet
        slide.notes_slide.notes_text_frame.text = note

    table_slide = presentation.slides.add_slide(title_only_layout)
    table_slide.shapes.title.text = "Master Theorem Cases"
    table = table_slide.shapes.add_table(
        4,
        3,
        Inches(0.9),
        Inches(1.8),
        Inches(8.2),
        Inches(2.6),
    ).table
    for column, header in enumerate(("Case", "Condition", "Result")):
        table.cell(0, column).text = header
    for row, values in enumerate(_MASTER_THEOREM_ROWS, start=1):
        for column, value in enumerate(values):
            table.cell(row, column).text = value
    table_slide.notes_slide.notes_text_frame.text = (
        "Work through case 2 with merge sort before assigning the problem set."
    )

    presentation.save(path)


def _write_markdown(path: Path) -> None:
    path.write_text(_MARKDOWN_NOTES, encoding="utf-8")


_PDF_PAGES: tuple[tuple[str, str], ...] = (
    (
        "Single-Source Shortest Paths",
        "A shortest-path algorithm takes a weighted directed graph and a source "
        "vertex, then reports the minimum total weight of any path from that "
        "source to every other vertex. Dijkstra's algorithm solves this problem "
        "in O((V + E) log V) time with a binary heap, but it depends on a "
        "critical assumption: every edge weight must be non-negative.\n\n"
        "The reason is the greedy invariant. Dijkstra removes the closest "
        "unfinished vertex from the priority queue and treats its distance as "
        "final. If a negative edge existed later in the graph, a path through "
        "an apparently distant vertex could still reduce the total, and the "
        "finalized distance would be wrong. Dijkstra never revisits a settled "
        "vertex, so it cannot recover from that mistake.\n\n"
        "Use Dijkstra when weights represent physical quantities that cannot be "
        "negative, such as distance, latency, or capacity. Reach for "
        "Bellman-Ford when weights can represent gains and losses, such as "
        "currency conversion or net energy change.",
    ),
    (
        "Bellman-Ford And Negative Edge Weights",
        "Bellman-Ford relaxes every edge in the graph V - 1 times. After "
        "iteration k, the algorithm has correctly computed every shortest path "
        "that uses at most k edges, so after V - 1 iterations it has computed "
        "every simple shortest path. The running time is O(V * E), which is "
        "slower than Dijkstra, and the reward for that cost is correctness in "
        "the presence of negative edge weights.\n\n"
        "The algorithm also detects negative cycles. Run one additional "
        "relaxation pass after the V - 1 required passes. If any edge can still "
        "be relaxed, some path reachable from the source enters a cycle whose "
        "total weight is negative, and no finite shortest path exists, because "
        "each additional trip around the cycle lowers the total further.\n\n"
        "Two practical refinements matter. First, stop early when a full pass "
        "relaxes no edge, because the distances have converged. Second, keep a "
        "predecessor pointer for each vertex so the path itself can be "
        "reconstructed, not only its weight.",
    ),
    (
        "All-Pairs Shortest Paths",
        "Floyd-Warshall computes shortest paths between every pair of vertices "
        "using dynamic programming over the set of permitted intermediate "
        "vertices. The recurrence considers whether routing a path from i to j "
        "through vertex k is shorter than the best route found so far, which "
        "yields a triple nested loop and O(V^3) running time with O(V^2) "
        "space.\n\n"
        "Loop order is not a stylistic choice. The intermediate vertex k must "
        "be the outermost loop, because the recurrence for stage k depends on "
        "results that already permit intermediate vertices 1 through k - 1. "
        "Placing k innermost produces values that are silently wrong on many "
        "graphs while still looking plausible on small examples.\n\n"
        "Negative edges are permitted. A negative cycle appears as a negative "
        "value on the diagonal of the distance matrix, since that would mean a "
        "vertex can reach itself at a total cost below zero. On sparse graphs, "
        "running Dijkstra from every vertex after Johnson's reweighting beats "
        "Floyd-Warshall.",
    ),
)


_PPTX_SLIDES: tuple[tuple[str, tuple[str, ...], str], ...] = (
    (
        "Asymptotic Notation",
        (
            "Big-O gives an upper bound: f(n) is O(g(n)) when f grows no faster "
            "than g beyond some threshold.",
            "Big-Omega gives a lower bound, and Big-Theta gives a tight bound "
            "when the upper and lower bounds match.",
            "Constant factors and lower-order terms are dropped, so 3n^2 + 40n "
            "and n^2 belong to the same class.",
            "Asymptotic class is not a promise about small inputs; insertion "
            "sort beats merge sort on short arrays.",
        ),
        "Ask the class why an O(n log n) algorithm can lose to an O(n^2) one in "
        "practice. Answer: constant factors and cache behaviour.",
    ),
    (
        "Dynamic Programming And Memoization",
        (
            "A problem suits dynamic programming when it has overlapping "
            "subproblems and an optimal substructure.",
            "Memoization is the top-down form: recurse as usual, but cache each "
            "subproblem result the first time it is computed.",
            "Tabulation is the bottom-up form: fill a table in dependency "
            "order, which avoids recursion depth limits.",
            "Naive Fibonacci recursion costs exponential time; memoization "
            "reduces it to linear time and linear space.",
        ),
        "Use Fibonacci as the opening example, then move to edit distance so "
        "students see a two-dimensional table.",
    ),
    (
        "Amortized Analysis",
        (
            "Amortized cost is the average cost per operation across a worst-case "
            "sequence, not an average over random inputs.",
            "A dynamic array that doubles its capacity has O(1) amortized append "
            "cost, even though a single resize costs O(n).",
            "The accounting method assigns credit to cheap operations to pay for "
            "later expensive ones.",
            "Amortized O(1) is not worst-case O(1); a real-time system may still "
            "need the worst-case guarantee.",
        ),
        "Contrast amortized with average-case explicitly. Students conflate them "
        "on the midterm every year.",
    ),
)


_MASTER_THEOREM_ROWS: tuple[tuple[str, str, str], ...] = (
    ("1", "f(n) grows slower than n^log_b(a)", "T(n) = Theta(n^log_b(a))"),
    ("2", "f(n) matches n^log_b(a)", "T(n) = Theta(n^log_b(a) * log n)"),
    ("3", "f(n) grows faster and is regular", "T(n) = Theta(f(n))"),
)


_MARKDOWN_NOTES = """# Exam Two Review Notes

## Choosing A Shortest-Path Algorithm

Pick the algorithm from the constraints, not from habit.

- Non-negative weights, one source: Dijkstra, O((V + E) log V).
- Negative weights possible, one source: Bellman-Ford, O(V * E).
- Every pair of vertices, dense graph: Floyd-Warshall, O(V^3).
- Every pair of vertices, sparse graph: Johnson's reweighting, then Dijkstra
  from each vertex.

## Recurrences

The master theorem applies to recurrences of the form T(n) = a * T(n / b) +
f(n). It does not apply when the subproblem sizes differ, as in T(n) =
T(n / 3) + T(2n / 3) + n, where a recursion tree gives the answer instead.

## Common Mistakes From Last Term

1. Running Dijkstra on a graph with negative edges and trusting the result.
2. Writing the Floyd-Warshall loops with the intermediate vertex innermost.
3. Reporting amortized cost as if it were a worst-case guarantee.
4. Forgetting the extra Bellman-Ford pass that detects negative cycles.
"""


__all__ = [
    "DEMO_COURSE_NAME",
    "DEMO_CORPUS_VERSION",
    "DEMO_MARKER_NAME",
    "DEMO_QUERY",
    "demo_source_root",
    "generate_demo_corpus",
    "run_demo",
]
