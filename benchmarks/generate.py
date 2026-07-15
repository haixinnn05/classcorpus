from __future__ import annotations

import argparse
from io import BytesIO
from pathlib import Path

import fitz
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches


def generate_corpus(output_dir: Path) -> dict[str, Path]:
    output_dir.mkdir(parents=True, exist_ok=True)
    generated = {
        "handout.pdf": make_pdf_fixture(output_dir / "handout.pdf"),
        "Lecture08.pptx": make_pptx_fixture(
            output_dir / "Lecture08.pptx",
            include_audit_slides=True,
        ),
    }
    generated["programming-glossary.md"] = _write_text_fixture(
        output_dir / "programming-glossary.md",
        "# Programming Glossary\n\n"
        + ("Dynamic programming language terminology. " * 30),
    )
    generated["edge-notes.txt"] = _write_text_fixture(
        output_dir / "edge-notes.txt",
        "Edge Notes\n" + ("Negative weights need careful handling. " * 30),
    )
    return generated


def make_pdf_fixture(path: Path) -> Path:
    document = fitz.open()
    page = document.new_page(width=612, height=8_000)
    content_block = ("precise-content " * 2_500).rstrip()
    text_blocks = [
        "Long lecture\n\n" + content_block,
        content_block,
        content_block,
        content_block,
    ]
    for index, text in enumerate(text_blocks):
        top = 30 + index * 1_980
        page.insert_textbox(
            fitz.Rect(30, top, 580, top + 1_950),
            text,
            fontsize=6,
        )

    image_page = document.new_page(width=612, height=792)
    image_page.insert_text((72, 96), "Diagram", fontsize=24)
    image_page.insert_text(
        (72, 128),
        "Bellman-Ford handles negative edges.",
        fontsize=12,
    )
    image_page.insert_image(
        fitz.Rect(72, 144, 540, 612),
        stream=png_bytes(),
    )

    document.save(path)
    document.close()
    return path


def make_pptx_fixture(path: Path, *, include_audit_slides: bool = False) -> Path:
    presentation = Presentation()

    slide_1 = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide_1.shapes.title.text = "Recurrences"
    textbox_1 = slide_1.shapes.add_textbox(
        Inches(1),
        Inches(1.5),
        Inches(5),
        Inches(1),
    )
    textbox_1.text_frame.text = "Divide-and-conquer recurrence trees."
    group = slide_1.shapes.add_group_shape()
    nested_group = group.shapes.add_group_shape()
    nested_textbox = nested_group.shapes.add_textbox(
        Inches(1),
        Inches(2.5),
        Inches(5),
        Inches(1),
    )
    nested_textbox.text_frame.text = "Nested group detail."
    picture_bytes = png_bytes()
    slide_1.shapes.add_picture(
        BytesIO(picture_bytes),
        Inches(7),
        Inches(1.5),
        Inches(1),
        Inches(1),
    )
    slide_1.shapes.add_picture(
        BytesIO(picture_bytes),
        Inches(7),
        Inches(3),
        Inches(1),
        Inches(1),
    )
    group.shapes.add_picture(
        BytesIO(picture_bytes),
        Inches(7),
        Inches(4.5),
        Inches(1),
        Inches(1),
    )
    slide_1.notes_slide.notes_text_frame.text = "Exact instructor note."

    extension_list = OxmlElement("p:extLst")
    extension = OxmlElement("p:ext")
    extension.set("uri", "classcorpus-audit-fixture")
    fallback_text = OxmlElement("a:t")
    fallback_text.text = "OOXML fallback detail"
    extension.append(fallback_text)
    extension_list.append(extension)
    textbox_1.element.append(extension_list)

    slide_2 = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide_2.shapes.title.text = "Dynamic Programming"
    textbox_2 = slide_2.shapes.add_textbox(
        Inches(1),
        Inches(1.5),
        Inches(5),
        Inches(1),
    )
    textbox_2.text_frame.text = "Memoization avoids repeated subproblems."
    table = slide_2.shapes.add_table(
        2,
        2,
        Inches(1),
        Inches(2.4),
        Inches(5),
        Inches(1.2),
    ).table
    table.cell(0, 0).text = "Problem"
    table.cell(0, 1).text = "State"
    table.cell(1, 0).text = "Fibonacci"
    table.cell(1, 1).text = "n"
    mapped_group = slide_2.shapes.add_group_shape()
    mapped_nested_group = mapped_group.shapes.add_group_shape()
    mapped_textbox = mapped_nested_group.shapes.add_textbox(
        Inches(1),
        Inches(4),
        Inches(5),
        Inches(1),
    )
    mapped_textbox.text_frame.text = "Mapped nested group detail."
    slide_2.notes_slide.notes_text_frame.text = "Use Fibonacci as the example."

    if include_audit_slides:
        _add_audit_slides(presentation)

    presentation.save(path)
    return path


def _add_audit_slides(presentation: Presentation) -> None:
    chart_slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    chart_slide.shapes.title.text = "Chart audit"
    chart_data = ChartData()
    chart_data.categories = ["A", "B"]
    chart_data.add_series("Series", (1, 2))
    chart_slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(1),
        Inches(1.5),
        Inches(5),
        Inches(3),
        chart_data,
    )

    equation_slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    equation_slide.shapes.title.text = "Equation audit"
    _append_extension(equation_slide, "m:oMath", "equation")

    smartart_slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    smartart_slide.shapes.title.text = "SmartArt audit"
    _append_extension(smartart_slide, "dgm:relIds", "smartart")

    ole_slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    ole_slide.shapes.title.text = "OLE audit"
    _append_extension(ole_slide, "p:oleObj", "ole")


def _append_extension(slide, child_tag: str, label: str) -> None:
    extension_list = OxmlElement("p:extLst")
    extension = OxmlElement("p:ext")
    extension.set("uri", f"classcorpus-{label}-fixture")
    if child_tag == "dgm:relIds":
        child = extension.makeelement(
            "{http://schemas.openxmlformats.org/drawingml/2006/diagram}relIds"
        )
    else:
        child = OxmlElement(child_tag)
    extension.append(child)
    extension_list.append(extension)
    slide.element.append(extension_list)


def png_bytes() -> bytes:
    pixmap = fitz.Pixmap(fitz.csRGB, fitz.IRect(0, 0, 32, 32), False)
    pixmap.clear_with(180)
    return pixmap.tobytes("png")


def _write_text_fixture(path: Path, content: str) -> Path:
    path.write_text(content, encoding="utf-8")
    return path


def main() -> int:
    parser = argparse.ArgumentParser(
        description="Generate the redistributable ClassCorpus benchmark corpus."
    )
    parser.add_argument("--output", type=Path, required=True)
    arguments = parser.parse_args()
    generated = generate_corpus(arguments.output.resolve())
    for name, path in generated.items():
        print(f"{name}: {path}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
