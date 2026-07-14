from io import BytesIO
from pathlib import Path

import fitz
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches


def make_pdf_fixture(path: Path) -> Path:
    document = fitz.open()
    page = document.new_page(width=612, height=8_000)
    content_block = ("precise-content " * 2_500).rstrip()
    text_blocks = [
        "Long lecture\n\n" + content_block,
        content_block,
        content_block,
        content_block,
    ]
    for index, text in enumerate(text_blocks):
        top = 30 + index * 1_980
        page.insert_textbox(
            fitz.Rect(30, top, 580, top + 1_950),
            text,
            fontsize=6,
        )

    image_page = document.new_page(width=612, height=792)
    image_page.insert_text((72, 96), "Diagram", fontsize=24)
    image_page.insert_text(
        (72, 128),
        "Bellman-Ford handles negative edges.",
        fontsize=12,
    )
    image_page.insert_image(
        fitz.Rect(72, 144, 540, 612),
        stream=_png_bytes(),
    )

    document.save(path)
    document.close()
    return path


def make_pptx_fixture(path: Path, *, include_audit_slides: bool = False) -> Path:
    presentation = Presentation()

    slide_1 = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide_1.shapes.title.text = "Recurrences"
    textbox_1 = slide_1.shapes.add_textbox(Inches(1), Inches(1.5), Inches(5), Inches(1))
    textbox_1.text_frame.text = "Divide-and-conquer recurrence trees."
    group = slide_1.shapes.add_group_shape()
    nested_group = group.shapes.add_group_shape()
    nested_textbox = nested_group.shapes.add_textbox(
        Inches(1),
        Inches(2.5),
        Inches(5),
        Inches(1),
    )
    nested_textbox.text_frame.text = "Nested group detail."
    slide_1.shapes.add_picture(
        BytesIO(_png_bytes()),
        Inches(7),
        Inches(1.5),
        Inches(1),
        Inches(1),
    )
    slide_1.notes_slide.notes_text_frame.text = "Exact instructor note."

    extension_list = OxmlElement("p:extLst")
    extension = OxmlElement("p:ext")
    extension.set("uri", "classcorpus-audit-fixture")
    fallback_text = OxmlElement("a:t")
    fallback_text.text = "OOXML fallback detail"
    extension.append(fallback_text)
    extension_list.append(extension)
    textbox_1.element.append(extension_list)

    slide_2 = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide_2.shapes.title.text = "Dynamic Programming"
    textbox_2 = slide_2.shapes.add_textbox(Inches(1), Inches(1.5), Inches(5), Inches(1))
    textbox_2.text_frame.text = "Memoization avoids repeated subproblems."
    table = slide_2.shapes.add_table(2, 2, Inches(1), Inches(2.4), Inches(5), Inches(1.2)).table
    table.cell(0, 0).text = "Problem"
    table.cell(0, 1).text = "State"
    table.cell(1, 0).text = "Fibonacci"
    table.cell(1, 1).text = "n"
    mapped_group = slide_2.shapes.add_group_shape()
    mapped_nested_group = mapped_group.shapes.add_group_shape()
    mapped_textbox = mapped_nested_group.shapes.add_textbox(
        Inches(1),
        Inches(4),
        Inches(5),
        Inches(1),
    )
    mapped_textbox.text_frame.text = "Mapped nested group detail."
    slide_2.notes_slide.notes_text_frame.text = "Use Fibonacci as the example."

    if include_audit_slides:
        chart_slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        chart_slide.shapes.title.text = "Chart audit"
        chart_data = ChartData()
        chart_data.categories = ["A", "B"]
        chart_data.add_series("Series", (1, 2))
        chart_slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            Inches(1),
            Inches(1.5),
            Inches(5),
            Inches(3),
            chart_data,
        )

        for title, tag in (
            ("Equation audit", "m:oMath"),
            ("OLE audit", "p:oleObj"),
        ):
            slide = presentation.slides.add_slide(presentation.slide_layouts[5])
            slide.shapes.title.text = title
            extension_list = OxmlElement("p:extLst")
            extension = OxmlElement("p:ext")
            extension.set("uri", f"classcorpus-{title.casefold().replace(' ', '-')}")
            extension.append(OxmlElement(tag))
            extension_list.append(extension)
            slide.element.append(extension_list)

    presentation.save(path)
    return path


def _png_bytes() -> bytes:
    pixmap = fitz.Pixmap(fitz.csRGB, fitz.IRect(0, 0, 32, 32), False)
    pixmap.clear_with(180)
    return pixmap.tobytes("png")
