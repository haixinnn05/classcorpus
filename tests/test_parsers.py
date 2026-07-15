from pathlib import Path
import subprocess

import fitz
import pytest
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from classcorpus.parsers import UnsupportedFormatError, _pdf_audit, parse_source
from tests.fixtures.make_fixtures import (
    _png_bytes,
    make_pdf_fixture,
    make_pptx_fixture,
)


@pytest.fixture
def pdf_fixture(tmp_path: Path) -> Path:
    return make_pdf_fixture(tmp_path / "lecture.pdf")


@pytest.fixture
def pptx_fixture(tmp_path: Path) -> Path:
    return make_pptx_fixture(
        tmp_path / "lecture.pptx",
        include_audit_slides=True,
    )


def test_pdf_preserves_pages_and_renders(pdf_fixture: Path, tmp_path: Path):
    records = parse_source(pdf_fixture, tmp_path / "renders")

    with fitz.open(pdf_fixture) as document:
        source_page = document[0]
        generated_text = source_page.get_text("text", sort=True)
        generated_blocks = [
            str(block[4])
            for block in source_page.get_text("blocks")
            if len(block) > 4 and "precise-content" in str(block[4])
        ]
        assert len(records) == document.page_count
        assert [record.ordinal for record in records] == list(
            range(1, document.page_count + 1)
        )
        assert len(generated_blocks) == 4
        assert (
            sum(block.count("precise-content") for block in generated_blocks) == 10_000
        )
        assert generated_text.count("precise-content") == 10_000
        assert len(generated_text) > 100_000
        assert generated_text.endswith("precise-content")

    assert records[0].raw_text == generated_text
    assert records[0].title == "Long lecture"
    assert records[0].speaker_notes == ""
    assert records[0].native_text_chars == len(records[0].raw_text)
    assert records[0].extraction_status == "text-extracted"
    assert records[1].extraction_status == "review-needed"
    assert "low-native-text" in records[1].extraction_reasons
    assert "embedded-image" in records[1].extraction_reasons
    assert records[1].has_visual_content is True
    assert Path(records[0].render_path).is_file()
    assert Path(records[1].render_path).is_file()


def test_pptx_preserves_content_and_embedded_assets_without_subprocesses(
    pptx_fixture: Path,
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
):
    monkeypatch.setattr(
        subprocess,
        "run",
        lambda *args, **kwargs: pytest.fail("PPTX parsing invoked a subprocess"),
    )
    records = parse_source(pptx_fixture, tmp_path / "renders")

    assert len(records) == len(Presentation(pptx_fixture).slides)
    assert [record.ordinal for record in records] == list(range(1, len(records) + 1))
    assert records[0].title == "Recurrences"
    assert "Nested group detail." in records[0].body_text
    assert records[0].raw_text.count("OOXML fallback detail") == 1
    assert "OOXML fallback detail" in records[0].body_text
    assert "unmapped-ooxml-text" in records[0].extraction_reasons
    assert "embedded-image" in records[0].extraction_reasons
    assert records[0].extraction_status == "review-needed"
    assert records[0].speaker_notes == "Exact instructor note."
    assert records[0].native_text_chars == len(records[0].raw_text)
    assert records[0].has_visual_content is True
    assert records[1].title == "Dynamic Programming"
    assert "Memoization avoids repeated subproblems." in records[1].body_text
    assert "Fibonacci" in records[1].body_text
    assert "Mapped nested group detail." in records[1].body_text
    assert records[1].speaker_notes == "Use Fibonacci as the example."
    assert "unmapped-ooxml-text" not in records[1].extraction_reasons
    assert records[1].extraction_reasons == ()
    assert records[1].extraction_status == "text-extracted"
    assert records[2].title == "Chart audit"
    assert records[2].extraction_reasons == ("chart-or-diagram",)
    assert records[2].extraction_status == "review-needed"
    assert records[2].has_visual_content is True
    assert records[3].title == "Equation audit"
    assert records[3].extraction_reasons == ("equation-or-embedded-object",)
    assert records[3].extraction_status == "review-needed"
    assert records[3].has_visual_content is True
    assert records[4].title == "OLE audit"
    assert records[4].extraction_reasons == ("equation-or-embedded-object",)
    assert records[4].extraction_status == "review-needed"
    assert records[4].has_visual_content is True
    assert all(record.render_path is None for record in records)

    source_slide = Presentation(pptx_fixture).slides[0]
    source_pictures = list(_pictures_in_source_order(source_slide.shapes))
    assert len(records[0].visual_assets) == 3
    assert len({asset.path for asset in records[0].visual_assets}) == 1
    assert all(
        Path(asset.path).read_bytes() == _png_bytes()
        for asset in records[0].visual_assets
    )
    assert (
        records[0].visual_assets[0].left,
        records[0].visual_assets[0].top,
        records[0].visual_assets[0].width,
        records[0].visual_assets[0].height,
    ) == (
        source_pictures[0].left,
        source_pictures[0].top,
        source_pictures[0].width,
        source_pictures[0].height,
    )


def _pictures_in_source_order(shapes):
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _pictures_in_source_order(shape.shapes)
        elif shape.shape_type in {
            MSO_SHAPE_TYPE.PICTURE,
            MSO_SHAPE_TYPE.LINKED_PICTURE,
        }:
            yield shape


class _PdfAuditPage:
    def __init__(self, words: list[tuple], blocks: list[tuple]):
        self.words = words
        self.blocks = blocks

    def get_images(self, *, full: bool):
        assert full is True
        return []

    def get_drawings(self):
        return []

    def get_text(self, mode: str):
        if mode == "words":
            return self.words
        if mode == "blocks":
            return self.blocks
        raise AssertionError(f"Unexpected extraction mode: {mode}")


def test_pdf_audit_flags_only_divergent_native_extractor_tokens():
    raw_text = "Café paths\ncost 7"
    matching_words = [
        (0.0, 0.0, 20.0, 10.0, "Café", 0, 0, 0),
        (21.0, 0.0, 40.0, 10.0, "paths", 0, 0, 1),
        (0.0, 12.0, 20.0, 22.0, "cost", 0, 1, 0),
        (21.0, 12.0, 25.0, 22.0, "7", 0, 1, 1),
    ]
    matching_blocks = [
        (0.0, 0.0, 40.0, 22.0, "Café paths\ncost 7", 0, 0),
    ]

    _, matching_reasons = _pdf_audit(
        _PdfAuditPage(matching_words, matching_blocks),
        raw_text,
    )
    divergent_words = matching_words + [
        (26.0, 12.0, 50.0, 22.0, "hidden", 0, 1, 2),
    ]
    _, divergent_reasons = _pdf_audit(
        _PdfAuditPage(divergent_words, matching_blocks),
        raw_text,
    )

    assert "native-extractor-disagreement" not in matching_reasons
    assert divergent_reasons == ("native-extractor-disagreement",)


def test_parse_source_rejects_unsupported_formats(tmp_path: Path):
    source = tmp_path / "lecture.txt"
    source.write_text("not a lecture", encoding="utf-8")

    with pytest.raises(UnsupportedFormatError):
        parse_source(source, tmp_path / "renders")
